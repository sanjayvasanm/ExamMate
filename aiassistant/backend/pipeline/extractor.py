import fitz  # PyMuPDF
import pdfplumber
import os
import tempfile

_ocr_instance = None

def get_ocr():
    global _ocr_instance
    if _ocr_instance is None:
        try:
            from paddleocr import PaddleOCR
            print("[AI] Initializing PaddleOCR model...")
            _ocr_instance = PaddleOCR(use_angle_cls=True, lang='en', show_log=False)
        except Exception as e:
            print(f"[AI] OCR Initialization failed: {e}")
            return None
    return _ocr_instance

def extract_pptx(pptx_path):
    """Extract text from PPTX slides and notes."""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        text_blocks = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    slide_text.append(f"[Speaker Notes]: {notes}")
            full_text = "\n".join(slide_text).strip()
            if full_text:
                text_blocks.append({"type": "pptx", "content": full_text, "page": i + 1})
        return text_blocks
    except Exception as e:
        print(f"PPTX Extraction error: {e}")
        return []

def extract_docx(docx_path):
    """Extract text from Word (.docx) files."""
    try:
        from docx import Document
        doc = Document(docx_path)
        text_blocks = []
        for i, para in enumerate(doc.paragraphs):
            if para.text.strip():
                text_blocks.append({"type": "docx", "content": para.text.strip(), "page": i + 1})
        
        # Also handle tables in Word
        for i, table in enumerate(doc.tables):
            table_data = []
            for row in table.rows:
                table_data.append(" | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()]))
            table_text = "\n".join(table_data).strip()
            if table_text:
                text_blocks.append({"type": "table_docx", "content": table_text, "page": "table_" + str(i+1)})
        return text_blocks
    except Exception as e:
        print(f"DOCX Extraction error: {e}")
        return []

def extract_image(img_path):
    """Extract text from images using OCR."""
    ocr_active = get_ocr()
    if ocr_active is None:
        return []
    try:
        result = ocr_active.ocr(img_path, cls=True)
        text_blocks = []
        if result and result[0]:
            full_text = "\n".join([line[1][0] for line in result[0]])
            text_blocks.append({"type": "image_ocr", "content": full_text, "page": 1})
        return text_blocks
    except Exception as e:
        print(f"Image OCR error: {e}")
        return []

def extract_txt(txt_path):
    """Extract text from raw text files."""
    try:
        with open(txt_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read().strip()
            if content:
                return [{"type": "txt", "content": content, "page": 1}]
        return []
    except Exception as e:
        print(f"TXT Extraction error: {e}")
        return []

def extract_pdf(pdf_path):
    """Extracts text and tables from PDF using PyMuPDF and pdfplumber."""
    text_blocks = []
    try:
        doc = fitz.open(pdf_path)
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            page_text = page.get_text("blocks")
            
            ocr_active = get_ocr()
            if len(page.get_text("text").strip()) < 50 and ocr_active is not None:
                pix = page.get_pixmap()
                img_path = os.path.join(tempfile.gettempdir(), f"page_{page_num}.png")
                pix.save(img_path)
                result = ocr_active.ocr(img_path, cls=True)
                if result and result[0]:
                    ocr_text = "\n".join([line[1][0] for line in result[0]])
                    text_blocks.append({"type": "ocr_pdf", "content": ocr_text, "page": page_num})
                if os.path.exists(img_path):
                    os.remove(img_path)
            else:
                for b in page_text:
                    if b[4].strip():
                        text_blocks.append({"type": "text_pdf", "content": b[4].strip(), "page": page_num})
        doc.close()
    except Exception as e:
        print(f"PyMuPDF Extraction error: {e}")

    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                tables = page.extract_tables()
                for table in tables:
                    table_text = "\n".join([" | ".join([str(cell) for cell in row if cell]) for row in table if row])
                    if table_text.strip():
                        text_blocks.append({"type": "table_pdf", "content": table_text, "page": page_num})
    except Exception as e:
        print(f"pdfplumber Extraction error: {e}")
    return text_blocks

def extract_text_hybrid(file_path):
    """Universal router for all document types."""
    ext = file_path.split(".")[-1].lower()
    if ext in ["pptx", "ppt"]:
        return extract_pptx(file_path)
    elif ext == "docx":
        return extract_docx(file_path)
    elif ext in ["png", "jpg", "jpeg"]:
        return extract_image(file_path)
    elif ext == "txt":
        return extract_txt(file_path)
    elif ext == "pdf":
        return extract_pdf(file_path)
    return []

def chunk_content(text_blocks, chunk_size=500, overlap=50):
    chunks = []
    current_chunk = ""
    for block in text_blocks:
        content = block['content'].strip()
        if not content: continue
        if len(content) > chunk_size:
            if current_chunk:
                chunks.append(current_chunk.strip())
                current_chunk = ""
            for i in range(0, len(content), chunk_size - overlap):
                chunks.append(content[i : i + chunk_size].strip())
        else:
            if len(current_chunk) + len(content) < chunk_size:
                current_chunk += "\n" + content
            else:
                chunks.append(current_chunk.strip())
                current_chunk = content
    if current_chunk: chunks.append(current_chunk.strip())
    return [c for c in chunks if c.strip()]
